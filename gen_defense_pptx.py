from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

BG       = RGBColor(0x06, 0x0D, 0x1A)
SURFACE  = RGBColor(0x0B, 0x16, 0x28)
SURFACE2 = RGBColor(0x0F, 0x1E, 0x35)
BORDER   = RGBColor(0x1A, 0x2D, 0x4A)
TEXT     = RGBColor(0xE8, 0xF0, 0xFE)
MUTED    = RGBColor(0x7B, 0x9A, 0xBF)
BRAND    = RGBColor(0x00, 0xD4, 0xFF)
ACCENT   = RGBColor(0xFF, 0x3C, 0x6F)
SIDEBAR  = RGBColor(0x04, 0x0A, 0x14)
DANGER   = RGBColor(0xDC, 0x26, 0x26)
SUCCESS  = RGBColor(0x05, 0x96, 0x69)
WARNING  = RGBColor(0xD9, 0x77, 0x06)
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
VIOLET   = RGBColor(0xA7, 0x8B, 0xFA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL = 16

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def sl():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, c):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    sh.line.fill.background()
    return sh


def tx(s, text, l, t, w, h, size=13, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def ml(s, lines, l, t, w, h, size=12, color=TEXT, bold=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, b, c = item, bold, color
        elif len(item) == 3:
            text, b, c = item[0], item[1], item[2]
        else:
            text, b, c = item[0], bold, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = text; r.font.size = Pt(size)
        r.font.bold = b; r.font.color.rgb = c
        r.font.name = "Calibri"
    return tb


def hdr(s, title, sub=""):
    box(s, 0, 0, 0.45, 7.5, SIDEBAR)
    box(s, 0, 0, 0.45, 0.85, BRAND)
    box(s, 0.45, 0, 12.88, 0.07, BRAND)
    tx(s, title, 0.65, 0.14, 11.8, 0.65, size=26, bold=True, color=BRAND)
    if sub:
        tx(s, sub, 0.65, 0.76, 11.8, 0.35, size=13, color=MUTED, italic=True)
    box(s, 0.45, 1.1, 12.88, 0.04, BORDER)


def pn(s, n):
    tx(s, f"{n} / {TOTAL}", 12.0, 7.15, 1.2, 0.3,
       size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def b(text):
    return "  ▸  " + text


def er_entity(s, name, col, fields, x, y, w):
    """Draw an ER entity box. fields = [(field_name, 'PK'|'FK'|'')]"""
    HEADER_H = 0.42
    ROW_H    = 0.30
    h = HEADER_H + len(fields) * ROW_H + 0.08
    box(s, x, y, w, h, SURFACE)
    box(s, x, y, w, HEADER_H, col)
    tx(s, name, x + 0.08, y + 0.05, w - 0.16, HEADER_H - 0.08,
       size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, x, y + HEADER_H, w, 0.03, col)
    for i, (fname, ftype) in enumerate(fields):
        fy = y + HEADER_H + 0.05 + i * ROW_H
        if ftype == 'PK':
            fc, prefix, fb = WARNING, "PK  ", True
        elif ftype == 'FK':
            fc, prefix, fb = ACCENT,  "FK  ", False
        else:
            fc, prefix, fb = MUTED,   "      ", False
        tx(s, prefix + fname, x + 0.1, fy, w - 0.18, ROW_H - 0.02,
           size=9, bold=fb, color=fc)
    return h


def er_line(s, x1, y1, x2, y2):
    c = s.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = MUTED
    c.line.width = Pt(1.5)


def fshape(s, stype, x, y, w, h, bg, border):
    sh = s.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    sh.line.color.rgb = border; sh.line.width = Pt(1.5)
    return sh


# =============================================================
# SLIDE 1 — Cover Page
# =============================================================
s = sl()
box(s, 0, 0, 13.33, 7.5, SIDEBAR)
box(s, 0, 0, 13.33, 0.07, BRAND)
box(s, 0, 7.43, 13.33, 0.07, BRAND)
box(s, 0, 0, 0.55, 7.5, RGBColor(0x00, 0x90, 0xAA))

for cx, cy, sz, col in [
    (9.8, 0.2, 4.2, SURFACE),
    (10.2, 0.0, 2.4, SURFACE2),
    (0.6, 5.2, 3.2, SURFACE),
]:
    c = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

box(s, 1.5, 0.22, 10.3, 0.9, SURFACE)
box(s, 1.5, 0.22, 10.3, 0.06, INDIGO)
tx(s, "[University / Institution Name]",
   1.7, 0.3, 9.9, 0.38, size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
tx(s, "Faculty: [Your Faculty / Department]   |   Academic Year: 2025 - 2026",
   1.7, 0.68, 9.9, 0.34, size=11, color=MUTED, align=PP_ALIGN.CENTER)

tx(s, "\U0001f510  FileVault",
   0.8, 1.35, 11.7, 1.25, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 3.0, 2.75, 7.3, 0.05, BRAND)

tx(s, "Design and Implementation of a Secure Web-Based",
   0.8, 2.9, 11.7, 0.48, size=17, color=BRAND, align=PP_ALIGN.CENTER)
tx(s, "File Monitoring System with Real-Time Alerts",
   0.8, 3.36, 11.7, 0.46, size=17, color=BRAND, align=PP_ALIGN.CENTER)

box(s, 3.0, 3.92, 7.3, 0.05, ACCENT)

tx(s, "Final Year Project Defense",
   0.8, 4.08, 11.7, 0.38, size=13, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

box(s, 1.5, 4.62, 10.3, 1.62, SURFACE)
box(s, 1.5, 4.62, 10.3, 0.06, BRAND)
ml(s, [
    ("Author: Dahirou-Bachar", True, TEXT),
    ("Student ID: [Your Student ID]   |   Programme: [BSc / MSc Cybersecurity]", False, MUTED),
    ("Supervisor: [Supervisor Name]   |   Defense Date: [Insert Date]", False, MUTED),
], 1.7, 4.76, 9.9, 1.38, size=12)

tx(s, "Python / Flask 3.0.3   ·   PostgreSQL   ·   AES-256   ·   WebAuthn/FIDO2   ·   Railway",
   0.8, 6.5, 11.7, 0.38, size=10, color=MUTED, align=PP_ALIGN.CENTER)
pn(s, 1)


# =============================================================
# SLIDE 2 — Plan of Presentation  (11 sections)
# =============================================================
s = sl()
hdr(s, "Plan of Presentation", "Overview of what will be covered")

items = [
    ("01", "Motivational Example",    "Why this problem matters in the real world",               ACCENT),
    ("02", "Aim and Objectives",      "What this project set out to achieve",                     BRAND),
    ("03", "Research Question",       "The central question guiding the work",                    INDIGO),
    ("04", "Review of Similar Work",  "What already exists and what is missing",                  VIOLET),
    ("05", "Research Methodology",    "How the project was designed and built",                   WARNING),
    ("06", "Technology Stack",        "Tools and frameworks used to build FileVault",             BRAND),
    ("07", "Results and Discussion",  "What was built and how it performed",                      SUCCESS),
    ("08", "Limitations",             "Known constraints of the current system",                  DANGER),
    ("09", "Future Work",             "Planned improvements and extensions",                      WARNING),
    ("10", "Conclusion",              "Key findings, contributions and closing statement",         MUTED),
]
for i, (num, title, desc, col) in enumerate(items):
    tp = 1.22 + i * 0.52
    box(s, 0.65, tp, 0.54, 0.46, col)
    tx(s, num, 0.65, tp + 0.05, 0.54, 0.36, size=11, bold=True,
       color=BG, align=PP_ALIGN.CENTER)
    box(s, 1.22, tp, 11.66, 0.46, SURFACE)
    tx(s, title, 1.40, tp + 0.05, 3.8,  0.36, size=11, bold=True, color=col)
    tx(s, desc,  5.32, tp + 0.05, 7.45, 0.36, size=10, color=MUTED)
pn(s, 2)


# =============================================================
# SLIDE 3 — Motivational Example
# =============================================================
s = sl()
hdr(s, "Motivational Example", "Why does this problem matter?")

box(s, 0.65, 1.22, 5.9, 5.9, SURFACE)
box(s, 0.65, 1.22, 5.9, 0.07, ACCENT)
tx(s, "\U0001f6a8  SolarWinds Attack — 2020", 0.82, 1.3, 5.6, 0.48,
   size=13, bold=True, color=ACCENT)
ml(s, [
    b("Attackers silently modified files inside\nthe SolarWinds software build system"),
    b("Tampering went UNDETECTED for months"),
    b("Malicious files sent to 18,000+ organisations\nincluding US government agencies"),
    ("", False, MUTED),
    ("Damage caused:", True, WARNING),
    ("Estimated $100 billion+ in losses", False, TEXT),
    ("Attackers had full network access for 9 months", False, TEXT),
    ("Months needed to fully remove the attackers", False, TEXT),
    ("", False, MUTED),
    ("The core failure:", True, ACCENT),
    ("No system was checking whether\nthe files had been changed.", False, TEXT),
    ("", False, MUTED),
    ("This is exactly what FileVault solves.", True, BRAND),
], 0.82, 1.85, 5.6, 5.1, size=11)

box(s, 6.78, 1.22, 6.2, 5.9, SURFACE)
box(s, 6.78, 1.22, 6.2, 0.07, BRAND)
tx(s, "\U0001f3e5  Hospital Scenario", 6.95, 1.3, 5.9, 0.48,
   size=13, bold=True, color=BRAND)
ml(s, [
    b("Patient records stored in a shared\ncloud folder"),
    b("Disgruntled employee quietly modifies\na prescription file"),
    b("No alert. No hash check.\nChange goes undetected for weeks."),
    ("", False, MUTED),
    ("Possible consequences:", True, DANGER),
    ("Wrong medication dosage administered", False, TEXT),
    ("Wrong drug — could trigger allergic reaction", False, TEXT),
    ("In worst case: patient death", False, TEXT),
    ("", False, MUTED),
    ("With FileVault:", True, SUCCESS),
    b("Watchdog detects file change instantly"),
    b("SHA-256 hash no longer matches baseline"),
    b("Admin alerted on dashboard and by email"),
    b("File locked until admin approves download"),
], 6.95, 1.85, 5.9, 5.1, size=11)
pn(s, 3)


# =============================================================
# SLIDE 4 — Aim and Objectives
# =============================================================
s = sl()
hdr(s, "Aim and Objectives")

box(s, 0.65, 1.22, 12.03, 0.88, SURFACE)
box(s, 0.65, 1.22, 12.03, 0.07, BRAND)
tx(s, "AIM", 0.82, 1.3, 1.0, 0.36, size=11, bold=True, color=BRAND)
tx(s, "To design and implement a secure web-based file storage and monitoring system that combines "
      "encryption, mandatory access control, multi-factor authentication, and real-time file integrity "
      "monitoring into a single unified application accessible to small and resource-constrained organisations.",
   1.9, 1.3, 10.6, 0.72, size=12, color=TEXT)

objectives = [
    ("1", "Implement multi-layer authentication: password login, TOTP MFA, and WebAuthn/FIDO2 biometric.", BRAND),
    ("2", "Enforce Mandatory Access Control with five clearance levels (L1 TOP SECRET to L5 UNCLASSIFIED).", INDIGO),
    ("3", "Encrypt all files at rest using AES-256 Fernet with a unique HKDF-SHA256 key per file.",        VIOLET),
    ("4", "Build a File Integrity Monitoring engine: SHA-256 baseline, real-time watchdog, and scheduler.", ACCENT),
    ("5", "Implement an administrator approval workflow using time-limited one-time tokens.",                WARNING),
    ("6", "Deliver real-time security alerts via WebSocket and email notification.",                         SUCCESS),
    ("7", "Evaluate the system through functional, security, usability, and performance testing.",           MUTED),
]
tx(s, "OBJECTIVES", 0.65, 2.24, 2.5, 0.38, size=12, bold=True, color=BRAND)
for i, (num, text, col) in enumerate(objectives):
    tp = 2.68 + i * 0.64
    box(s, 0.65, tp, 0.48, 0.54, col)
    tx(s, num, 0.65, tp + 0.07, 0.48, 0.38, size=12, bold=True,
       color=BG, align=PP_ALIGN.CENTER)
    box(s, 1.16, tp, 11.52, 0.54, SURFACE if i % 2 == 0 else SURFACE2)
    tx(s, text, 1.3, tp + 0.09, 11.25, 0.38, size=11, color=TEXT)
pn(s, 4)


# =============================================================
# SLIDE 5 — Research Question
# =============================================================
s = sl()
hdr(s, "Research Question")

box(s, 0.65, 1.22, 12.03, 1.35, SURFACE)
box(s, 0.65, 1.22, 0.07, 1.35, BRAND)
tx(s, "MAIN QUESTION", 0.82, 1.3, 3.5, 0.38, size=11, bold=True, color=BRAND)
tx(s, '"How can encrypted file storage, mandatory access control, and real-time file integrity\n'
      'monitoring be integrated into a single lightweight web-based system accessible to\n'
      'small and resource-constrained organisations?"',
   0.82, 1.68, 11.7, 0.82, size=14, bold=True, color=WHITE, italic=True)

tx(s, "SUB-QUESTIONS", 0.65, 2.74, 4.0, 0.38, size=12, bold=True, color=BRAND)
subs = [
    ("Q1", "What authentication mechanisms provide the strongest protection against unauthorised access?",   BRAND),
    ("Q2", "How can file integrity be continuously monitored without significant performance impact?",        INDIGO),
    ("Q3", "What access control model best supports file classification at different sensitivity levels?",    VIOLET),
    ("Q4", "How can real-time alerts be delivered to administrators when a security event is detected?",      ACCENT),
    ("Q5", "Can enterprise-level file security be implemented affordably using open-source tools?",           SUCCESS),
]
for i, (num, text, col) in enumerate(subs):
    tp = 3.18 + i * 0.72
    box(s, 0.65, tp, 0.52, 0.6, col)
    tx(s, num, 0.65, tp + 0.1, 0.52, 0.42, size=11, bold=True,
       color=BG, align=PP_ALIGN.CENTER)
    box(s, 1.2, tp, 11.48, 0.6, SURFACE if i % 2 == 0 else SURFACE2)
    tx(s, text, 1.35, tp + 0.1, 11.2, 0.42, size=12, color=TEXT)
pn(s, 5)


# =============================================================
# SLIDE 6 — Review of Similar Projects
# =============================================================
s = sl()
hdr(s, "Review of Similar Projects", "What exists — and what is missing")

tools = [
    ("Tripwire",      ACCENT,  "FIM only — no file storage, no encryption, no web UI. Costs $10,000+/year."),
    ("Varonis",       VIOLET,  "File access monitoring — no encryption, no FIDO2 login, extremely expensive."),
    ("Dropbox Biz",   INDIGO,  "Cloud storage — no integrity hashing, no tamper detection, no clearance levels."),
    ("Nextcloud",     BRAND,   "Self-hosted storage — no SHA-256 baseline, no severity scoring, no real-time FIM."),
    ("OSSEC / Wazuh", WARNING, "Host IDS with FIM — no file vault, no web interface, requires complex agent setup."),
]
for i, (name, col, gap) in enumerate(tools):
    tp = 1.22 + i * 0.98
    box(s, 0.65, tp, 2.2, 0.84, col)
    tx(s, name, 0.65, tp + 0.19, 2.2, 0.45, size=13, bold=True,
       color=BG, align=PP_ALIGN.CENTER)
    box(s, 2.88, tp, 9.8, 0.84, SURFACE)
    tx(s, "✗  " + gap, 3.05, tp + 0.18, 9.55, 0.48, size=12, color=TEXT)

box(s, 0.65, 6.17, 12.03, 0.95, SURFACE)
box(s, 0.65, 6.17, 12.03, 0.06, SUCCESS)
tx(s, "✓  FileVault uniquely combines ALL FIVE pillars: encrypted storage  ·  mandatory access control  "
      "·  MFA  ·  real-time FIM  ·  WebAuthn  — in one free open-source web application.",
   0.82, 6.28, 11.75, 0.72, size=12, bold=True, color=SUCCESS)
pn(s, 6)


# =============================================================
# SLIDE 7 — Feature Comparison Table
# =============================================================
s = sl()
hdr(s, "Feature Comparison Table", "FileVault vs existing tools")

cols   = ["Feature", "Tripwire", "Varonis", "Dropbox", "Nextcloud", "FileVault"]
widths = [3.8, 1.5, 1.5, 1.5, 1.65, 1.85]
lefts  = [0.65, 4.48, 5.98, 7.48, 8.98, 10.63]

box(s, 0.65, 1.22, 12.03, 0.52, SURFACE2)
for j, (col, lft, w) in enumerate(zip(cols, lefts, widths)):
    c = BRAND if j == 0 else (SUCCESS if j == 5 else MUTED)
    tx(s, col, lft + 0.08, 1.28, w - 0.1, 0.38,
       size=11, bold=True, color=c,
       align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rows = [
    ("AES-256 File Encryption",   ["No", "No", "Partial", "No",  "✓"]),
    ("Per-file Key Derivation",   ["No", "No", "No",      "No",  "✓"]),
    ("Real-time FIM + Alerts",    ["✓",  "✓",  "No",      "No",  "✓"]),
    ("Clearance-Level RBAC",      ["No", "No", "No",      "No",  "✓"]),
    ("Admin Approval Workflow",   ["No", "No", "No",      "No",  "✓"]),
    ("WebAuthn / FIDO2 Login",    ["No", "No", "No",      "No",  "✓"]),
    ("TOTP MFA",                  ["No", "No", "✓",       "✓",   "✓"]),
    ("Free and Open Source",      ["No", "No", "No",      "✓",   "✓"]),
]
for i, (feat, vals) in enumerate(rows):
    tp = 1.8 + i * 0.62
    bg = SURFACE if i % 2 == 0 else SURFACE2
    box(s, 0.65, tp, 12.03, 0.58, bg)
    tx(s, feat, lefts[0] + 0.08, tp + 0.1, widths[0] - 0.1, 0.38, size=11, color=TEXT)
    for j, val in enumerate(vals):
        col = SUCCESS if val == "✓" else (MUTED if val == "No" else WARNING)
        tx(s, val, lefts[j + 1] + 0.08, tp + 0.1, widths[j + 1] - 0.1, 0.38,
           size=12, bold=(val == "✓"), color=col, align=PP_ALIGN.CENTER)
pn(s, 7)


# =============================================================
# SLIDE 8 — Research Methodology
# =============================================================
s = sl()
hdr(s, "Research Methodology", "Design Science Research + Scrum Agile")

box(s, 0.65, 1.22, 5.85, 5.9, SURFACE)
box(s, 0.65, 1.22, 5.85, 0.07, INDIGO)
tx(s, "Design Science Research (DSR)", 0.82, 1.3, 5.55, 0.45,
   size=13, bold=True, color=INDIGO)
ml(s, [
    ("Used in Computer Science when the goal", False, TEXT),
    ("is to BUILD and EVALUATE an artefact", False, TEXT),
    ("that solves a real-world problem.", False, TEXT),
    ("", False, MUTED),
    ("Three core activities:", True, INDIGO),
    b("Problem Identification"),
    ("    Gap in existing file security tools", False, MUTED),
    b("Design and Development"),
    ("    The four Scrum sprints", False, MUTED),
    b("Evaluation"),
    ("    Functional, security, usability,", False, MUTED),
    ("    and performance testing", False, MUTED),
    ("", False, MUTED),
    ("FileVault IS the artefact.", True, BRAND),
    ("All roles handled by one person.", False, MUTED),
], 0.82, 1.82, 5.55, 5.1, size=12)

box(s, 6.73, 1.22, 6.0, 5.9, SURFACE)
box(s, 6.73, 1.22, 6.0, 0.07, BRAND)
tx(s, "Scrum Framework  (Solo Adaptation)", 6.9, 1.3, 5.7, 0.45,
   size=13, bold=True, color=BRAND)

sprints = [
    ("Sprint 1", "Authentication",       "Registration, login, TOTP MFA,\nWebAuthn/FIDO2, rate limiting",   BRAND),
    ("Sprint 2", "File Vault & RBAC",    "Upload, AES-256 encryption, HKDF,\n5-level clearance, quota",     INDIGO),
    ("Sprint 3", "FIM & Alerts",         "SHA-256 baseline, watchdog, scheduler,\n8-factor scoring, alerts", ACCENT),
    ("Sprint 4", "Testing & Deployment", "Functional, security, usability tests,\nRailway cloud deployment", SUCCESS),
]
for i, (sp, title, desc, col) in enumerate(sprints):
    tp = 1.82 + i * 1.3
    box(s, 6.73, tp, 0.9, 1.18, col)
    tx(s, sp, 6.73, tp + 0.35, 0.9, 0.45,
       size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    box(s, 7.66, tp, 4.98, 1.18, SURFACE2)
    tx(s, title, 7.8, tp + 0.08, 4.7, 0.38, size=12, bold=True, color=col)
    tx(s, desc,  7.8, tp + 0.52, 4.7, 0.58, size=11, color=MUTED)
pn(s, 8)


# =============================================================
# SLIDE 9 — System Architecture (File Access Flow)
# =============================================================
s = sl()
hdr(s, "System Architecture — File Access Flow",
    "How FileVault controls access to stored files")

# ── layout constants ─────────────────────────────────────────
CX     = 4.5
OW, OH = 2.4, 0.52
DW, DH = 3.8, 0.76
BW, BH = 3.1, 0.52
OX     = CX - OW / 2      # 3.3
DX     = CX - DW / 2      # 2.6
DENY_X = 9.5

Y_START = 1.30
Y_D1    = 2.20
Y_D2    = 3.38
Y_D3    = 4.56
Y_GRANT = 5.68

# ── main column shapes ────────────────────────────────────────
fshape(s, 9, OX, Y_START, OW, OH, BRAND, BRAND)
tx(s, "START", OX, Y_START + 0.12, OW, 0.28,
   size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

fshape(s, 4, DX, Y_D1, DW, DH, SURFACE, BRAND)
tx(s, "User Authenticated?", CX - 1.3, Y_D1 + 0.24, 2.6, 0.28,
   size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

fshape(s, 4, DX, Y_D2, DW, DH, SURFACE, INDIGO)
tx(s, "Clearance Level OK?", CX - 1.3, Y_D2 + 0.24, 2.6, 0.28,
   size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

fshape(s, 4, DX, Y_D3, DW, DH, SURFACE, VIOLET)
tx(s, "Admin Approved?", CX - 1.3, Y_D3 + 0.24, 2.6, 0.28,
   size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

fshape(s, 9, OX, Y_GRANT, OW, OH, SUCCESS, SUCCESS)
tx(s, "GRANT ACCESS", OX, Y_GRANT + 0.12, OW, 0.28,
   size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

# ── diamond centres (for connectors and DENY boxes) ───────────
d1_cy = Y_D1 + DH / 2
d2_cy = Y_D2 + DH / 2
d3_cy = Y_D3 + DH / 2

# ── DENY outcome boxes ────────────────────────────────────────
fshape(s, 1, DENY_X, d1_cy - BH / 2, BW, BH, DANGER, DANGER)
tx(s, "DENY: Not Authenticated", DENY_X + 0.1, d1_cy - 0.14, BW - 0.2, 0.28,
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

fshape(s, 1, DENY_X, d2_cy - BH / 2, BW, BH, DANGER, DANGER)
tx(s, "DENY: Insufficient Clearance", DENY_X + 0.1, d2_cy - 0.14, BW - 0.2, 0.28,
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

fshape(s, 1, DENY_X, d3_cy - BH / 2, BW, BH, DANGER, DANGER)
tx(s, "DENY: Not Approved", DENY_X + 0.1, d3_cy - 0.14, BW - 0.2, 0.28,
   size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── connector lines ───────────────────────────────────────────
er_line(s, CX, Y_START + OH,  CX, Y_D1)
er_line(s, CX, Y_D1 + DH,     CX, Y_D2)
er_line(s, CX, Y_D2 + DH,     CX, Y_D3)
er_line(s, CX, Y_D3 + DH,     CX, Y_GRANT)
er_line(s, DX + DW, d1_cy,    DENY_X, d1_cy)
er_line(s, DX + DW, d2_cy,    DENY_X, d2_cy)
er_line(s, DX + DW, d3_cy,    DENY_X, d3_cy)

# ── Yes / No labels ───────────────────────────────────────────
tx(s, "Yes", CX + 0.12, Y_D1 + DH + 0.04, 0.6, 0.24, size=9, bold=True, color=SUCCESS)
tx(s, "Yes", CX + 0.12, Y_D2 + DH + 0.04, 0.6, 0.24, size=9, bold=True, color=SUCCESS)
tx(s, "Yes", CX + 0.12, Y_D3 + DH + 0.04, 0.6, 0.24, size=9, bold=True, color=SUCCESS)
tx(s, "No", DX + DW + 0.12, d1_cy - 0.28, 0.55, 0.24, size=9, bold=True, color=DANGER)
tx(s, "No", DX + DW + 0.12, d2_cy - 0.28, 0.55, 0.24, size=9, bold=True, color=DANGER)
tx(s, "No", DX + DW + 0.12, d3_cy - 0.28, 0.55, 0.24, size=9, bold=True, color=DANGER)

pn(s, 9)


# =============================================================
# SLIDE 10 — Iterative & Incremental Development
# =============================================================
s = sl()
hdr(s, "Iterative & Incremental Development",
    "How Scrum was applied across the four sprints")

box(s, 0.65, 1.22, 5.75, 5.9, SURFACE)
box(s, 0.65, 1.22, 5.75, 0.07, INDIGO)
tx(s, "What these terms mean", 0.82, 1.3, 5.45, 0.42,
   size=13, bold=True, color=INDIGO)

ml(s, [
    ("ITERATIVE", True, BRAND),
    ("Each sprint revisited and improved", False, TEXT),
    ("on work done in the previous sprint.", False, TEXT),
    ("Feedback was used to refine before", False, TEXT),
    ("moving on to the next phase.", False, TEXT),
    ("", False, MUTED),
    ("Example:", True, INDIGO),
    ("Sprint 2 improved Sprint 1 authentication", False, MUTED),
    ("before adding the file vault layer.", False, MUTED),
    ("Sprint 4 fixed all bugs found across", False, MUTED),
    ("sprints 1, 2, and 3.", False, MUTED),
    ("", False, MUTED),
    ("INCREMENTAL", True, BRAND),
    ("Each sprint delivered a working,", False, TEXT),
    ("testable piece of the system.", False, TEXT),
    ("No sprint left the system broken.", False, TEXT),
    ("Each built on top of the previous.", False, TEXT),
], 0.82, 1.82, 5.45, 5.0, size=12)

box(s, 6.65, 1.22, 6.03, 5.9, SURFACE)
box(s, 6.65, 1.22, 6.03, 0.07, BRAND)
tx(s, "How the system was built up", 6.82, 1.3, 5.73, 0.42,
   size=13, bold=True, color=BRAND)

sprint_rows = [
    ("Sprint 1", [("Authentication", BRAND)]),
    ("Sprint 2", [("Authentication", BRAND), ("File Vault + RBAC", INDIGO)]),
    ("Sprint 3", [("Authentication", BRAND), ("File Vault + RBAC", INDIGO), ("FIM + Alerts", ACCENT)]),
    ("Sprint 4", [("Authentication", BRAND), ("File Vault + RBAC", INDIGO), ("FIM + Alerts", ACCENT), ("Testing + Deploy", SUCCESS)]),
]

ROW_H = 1.15
for i, (sp_label, segments) in enumerate(sprint_rows):
    row_y = 1.88 + i * (ROW_H + 0.06)
    box(s, 6.65, row_y, 1.0, ROW_H, SURFACE2)
    tx(s, sp_label, 6.65, row_y + 0.35, 1.0, 0.42,
       size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    seg_w = 4.9 / 4
    for j, (seg_label, seg_col) in enumerate(segments):
        sx = 7.7 + j * (seg_w + 0.04)
        box(s, sx, row_y + 0.08, seg_w, ROW_H - 0.16, seg_col)
        tx(s, seg_label, sx + 0.05, row_y + 0.35, seg_w - 0.1, 0.42,
           size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)

tx(s, "Each sprint added a new layer — by Sprint 4 the system was complete, tested, and live.",
   6.82, 6.7, 5.73, 0.36, size=10, color=MUTED, italic=True)
pn(s, 10)


# =============================================================
# SLIDE 11 — Technology Stack
# =============================================================
s = sl()
hdr(s, "Technology Stack", "Tools and frameworks used to build FileVault")

tools_cards = [
    ("Backend", BRAND, [
        "Python 3.x",
        "Flask 3.0.3 web framework",
        "Gunicorn WSGI server",
    ]),
    ("Database", INDIGO, [
        "PostgreSQL (relational DB)",
        "SQLAlchemy 2.0 ORM",
        "Alembic migrations",
    ]),
    ("Encryption & Security", VIOLET, [
        "AES-256 Fernet (files at rest)",
        "HKDF-SHA256 per-file key",
        "PBKDF2 password hashing",
    ]),
    ("Authentication", ACCENT, [
        "Flask-Login (session management)",
        "pyotp — TOTP MFA",
        "webauthn 2.2.0 — FIDO2",
    ]),
    ("FIM & Real-Time Alerts", SUCCESS, [
        "Watchdog 4.0.1 (real-time FIM)",
        "Flask-SocketIO 5.3.6 (WebSocket)",
        "Resend API (email alerts)",
    ]),
    ("Deployment", WARNING, [
        "Railway cloud platform",
        "Gunicorn production server",
        "GitHub version control",
    ]),
]

CARD_W = 3.92; CARD_H = 2.57; GAP_X = 0.12; GAP_Y = 0.10
SX = 0.65;     SY = 1.22

for i, (title, col, bullets) in enumerate(tools_cards):
    col_i = i % 3; row_i = i // 3
    lft = SX + col_i * (CARD_W + GAP_X)
    tp  = SY + row_i * (CARD_H + GAP_Y)
    box(s, lft, tp, CARD_W, CARD_H, SURFACE)
    box(s, lft, tp, CARD_W, 0.07, col)
    box(s, lft, tp + 0.07, 0.05, CARD_H - 0.07, col)
    tx(s, title, lft + 0.18, tp + 0.14, CARD_W - 0.25, 0.42,
       size=13, bold=True, color=col)
    box(s, lft + 0.18, tp + 0.62, CARD_W - 0.35, 0.03, BORDER)
    for j, point in enumerate(bullets):
        tx(s, "▸  " + point,
           lft + 0.18, tp + 0.74 + j * 0.58,
           CARD_W - 0.28, 0.50, size=11, color=TEXT)
pn(s, 11)


# =============================================================
# SLIDE 12 — Results and Discussion
# =============================================================
s = sl()
hdr(s, "Results and Discussion", "All seven objectives met — system fully deployed")

tests = [
    ("Functional Testing",  SUCCESS, "All core user flows passed:\nregistration, login, MFA,\nupload, download, FIM, alerts."),
    ("Security Testing",    SUCCESS, "Encryption, hashing, lockout,\nrate limiting, token expiry,\nCSRF protection all verified."),
    ("Usability Testing",   SUCCESS, "System navigable without\ntechnical knowledge. Interface\nclear and responsive."),
    ("Performance Testing", SUCCESS, "Acceptable response times\nunder normal load. SocketIO\nalerts delivered in real time."),
]
for i, (title, col, desc) in enumerate(tests):
    lft = 0.65 + i * 3.17
    box(s, lft, 1.22, 2.9, 2.52, SURFACE)
    box(s, lft, 1.22, 2.9, 0.07, col)
    tx(s, "✓  " + title, lft + 0.15, 1.3, 2.65, 0.46, size=12, bold=True, color=col)
    tx(s, desc, lft + 0.15, 1.88, 2.65, 1.74, size=11, color=TEXT)

box(s, 0.65, 3.9, 12.03, 2.72, SURFACE)
box(s, 0.65, 3.9, 12.03, 0.07, BRAND)
tx(s, "DISCUSSION", 0.82, 3.98, 3.0, 0.38, size=12, bold=True, color=BRAND)
ml(s, [
    b("FileVault delivers all 7 objectives as a deployed, working web application"),
    b("Combines capabilities previously split across Tripwire, Varonis, and Dropbox into one free app"),
    b("The 8-factor FIM severity engine is an original contribution — no reviewed tool implements this"),
    b("WebAuthn/FIDO2 biometric login places this above typical student-level project work"),
    b("Proves enterprise-level file security is achievable with open-source tools at zero licensing cost"),
    b("All tests passed: functional, security, usability, and performance"),
], 0.82, 4.42, 11.7, 2.1, size=12)
pn(s, 12)


# =============================================================
# SLIDE 13 — Limitations
# =============================================================
s = sl()
hdr(s, "Limitations", "Known constraints of the current system")

lims = [
    ("In-Memory Rate Limiting", DANGER,
     "Rate limiting is stored in application memory and resets every time the server "
     "restarts or redeploys. In a production environment with multiple instances or "
     "frequent restarts, an attacker could bypass the rate limit simply by waiting "
     "for a restart. A Redis-backed solution would persist rate limiting across all "
     "restarts and scale correctly to multiple server instances."),
    ("No Encryption Key Rotation", WARNING,
     "The master encryption key used to derive per-file keys is static after "
     "deployment and never changes. If the key were ever compromised, all files "
     "encrypted under it would be at risk. A proper key rotation mechanism — "
     "combined with Hardware Security Module (HSM) integration — would allow "
     "periodic key changes without re-encrypting all files at once."),
    ("HTTPS Not Enforced by the Application", WARNING,
     "HTTPS is provided by Railway's reverse proxy layer, not enforced within "
     "Flask itself. If the application were accessed directly without the proxy "
     "(e.g., during local development or a misconfigured deployment), "
     "data would be transmitted unencrypted. Enforcing HTTPS at the application "
     "level using Flask-Talisman would close this gap."),
]

for i, (title, col, desc) in enumerate(lims):
    tp = 1.22 + i * 1.88
    box(s, 0.65, tp, 12.03, 1.76, SURFACE)
    box(s, 0.65, tp, 12.03, 0.07, col)
    box(s, 0.65, tp + 0.07, 0.07, 1.69, col)
    tx(s, title, 0.9, tp + 0.12, 5.0, 0.42, size=14, bold=True, color=col)
    tx(s, desc, 0.9, tp + 0.58, 11.5, 1.1, size=11, color=TEXT)
pn(s, 13)


# =============================================================
# SLIDE 14 — Future Work
# =============================================================
s = sl()
hdr(s, "Future Work", "Planned improvements and extensions to FileVault")

future = [
    ("Redis-backed rate limiting",            BRAND,
     "Replace in-memory rate limiting with Redis so limits persist across restarts and scale horizontally."),
    ("Encryption key rotation + HSM",         INDIGO,
     "Add scheduled key rotation and Hardware Security Module integration for enterprise-grade key management."),
    ("MIME type validation on upload",        VIOLET,
     "Validate actual file content type — not just the extension — to block disguised executable uploads."),
    ("SIEM / syslog integration",             WARNING,
     "Forward audit logs to external SIEM platforms (e.g. Splunk, Elastic) for centralised incident response."),
    ("Two-admin approval for baseline reset", ACCENT,
     "Require two administrators to sign off on critical baseline resets as a defence against insider threat."),
    ("Mobile push alert notifications",       SUCCESS,
     "Send real-time push notifications to administrator mobile devices on CRITICAL or HIGH severity events."),
]

for i, (title, col, desc) in enumerate(future):
    tp = 1.22 + i * 0.9
    box(s, 0.65, tp, 12.03, 0.82, SURFACE if i % 2 == 0 else SURFACE2)
    box(s, 0.65, tp, 0.07, 0.82, col)
    tx(s, "▸  " + title, 0.9, tp + 0.08, 4.2, 0.38, size=12, bold=True, color=col)
    tx(s, desc, 5.2, tp + 0.08, 7.3, 0.64, size=11, color=TEXT)
pn(s, 14)


# =============================================================
# SLIDE 15 — Conclusion
# =============================================================
s = sl()
box(s, 0, 0, 13.33, 7.5, SIDEBAR)
box(s, 0, 0, 13.33, 0.07, BRAND)
box(s, 0, 7.43, 13.33, 0.07, BRAND)

tx(s, "Conclusion", 0.8, 1.2, 11.7, 1.1,
   size=44, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
box(s, 2.8, 2.45, 7.73, 0.05, BRAND)

tx(s,
   "I would like to sincerely thank the President of the Jury "
   "and all honourable members for your time and attention. "
   "FileVault was conceived, designed, built, tested, and deployed "
   "entirely by one person. It has been a demanding but deeply rewarding journey, "
   "and I am now fully open to your questions.",
   1.5, 2.7, 10.3, 3.8,
   size=18, color=TEXT, align=PP_ALIGN.CENTER, italic=True)
pn(s, 15)


# =============================================================
# SLIDE 16 — Thank You / Q&A
# =============================================================
s = sl()
box(s, 0, 0, 13.33, 7.5, SIDEBAR)
box(s, 0, 0, 13.33, 0.07, BRAND)
box(s, 0, 7.43, 13.33, 0.07, BRAND)
box(s, 0, 0, 0.55, 7.5, RGBColor(0x00, 0x90, 0xAA))

for cx, cy, sz, col in [
    (9.5, 0.2, 4.5, SURFACE),
    (0.6, 4.8, 3.5, SURFACE),
]:
    c = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

tx(s, "Thank You", 0.8, 1.5, 11.7, 1.2,
   size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 3.0, 2.85, 7.3, 0.06, BRAND)
tx(s, "Questions & Discussion", 0.8, 3.05, 11.7, 0.58,
   size=22, color=BRAND, align=PP_ALIGN.CENTER)
box(s, 3.0, 3.82, 7.3, 0.06, ACCENT)

feat_items = [
    ("\U0001f512  AES-256 Encryption",    BRAND),
    ("\U0001f441  Real-Time FIM",          ACCENT),
    ("\U0001f510  WebAuthn / FIDO2",       VIOLET),
    ("\U0001f6e1  5-Level MAC",            INDIGO),
    ("\U0001f514  Live Alerts + Audit",    SUCCESS),
]
for i, (feat, col) in enumerate(feat_items):
    lft = 1.2 + i * 2.22
    box(s, lft, 4.1, 2.0, 0.62, SURFACE)
    box(s, lft, 4.1, 2.0, 0.06, col)
    tx(s, feat, lft + 0.1, 4.18, 1.85, 0.44, size=9, bold=True, color=col)

tx(s, "Author: Dahirou-Bachar   ·   FileVault   ·   Python / Flask / PostgreSQL / Railway",
   0.8, 5.1, 11.7, 0.38, size=11, color=MUTED, align=PP_ALIGN.CENTER)
tx(s, "Design and Implementation of a Secure Web-Based File Monitoring System with Real-Time Alerts",
   0.8, 5.58, 11.7, 0.38, size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
pn(s, 16)


# =============================================================
out = r"c:\Users\user\Desktop\FileVault_Defense_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
