from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── FileVault exact colour palette (from static/css) ────────────────
BG        = RGBColor(0x06, 0x0D, 0x1A)   # --bg
SURFACE   = RGBColor(0x0B, 0x16, 0x28)   # --surface
SURFACE2  = RGBColor(0x0F, 0x1E, 0x35)   # --surface-2
BORDER    = RGBColor(0x1A, 0x2D, 0x4A)   # --border
TEXT      = RGBColor(0xE8, 0xF0, 0xFE)   # --text
MUTED     = RGBColor(0x7B, 0x9A, 0xBF)   # --text-muted
BRAND     = RGBColor(0x00, 0xD4, 0xFF)   # --brand  (cyan)
ACCENT    = RGBColor(0xFF, 0x3C, 0x6F)   # --accent (hot-pink)
SIDEBAR   = RGBColor(0x04, 0x0A, 0x14)   # --sidebar-bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DANGER    = RGBColor(0xDC, 0x26, 0x26)   # --danger
SUCCESS   = RGBColor(0x05, 0x96, 0x69)   # --success
WARNING   = RGBColor(0xD9, 0x77, 0x06)   # --warning
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)   # --brand (light theme)
VIOLET    = RGBColor(0xA7, 0x8B, 0xFA)   # soft purple

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TOTAL   = 16

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def rect(s, l, t, w, h, color):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(s, text, l, t, w, h, size=13, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return tb


def mtxt(s, lines, l, t, w, h, size=12, color=TEXT, bold=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, b, c = item, bold, color
        elif len(item) == 3:
            text, b, c = item
        elif len(item) == 2:
            text, b, c = item[0], item[1], color
        else:
            text, b, c = item[0], bold, color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = "Calibri"
    return tb


def header(s, title, sub=""):
    # sidebar strip
    rect(s, 0, 0, 0.5, 7.5, SIDEBAR)
    rect(s, 0, 0, 0.5, 0.9, BRAND)
    # top bar
    rect(s, 0.5, 0, 12.83, 0.08, BRAND)
    txt(s, title, 0.7, 0.15, 11.5, 0.65,
        size=26, bold=True, color=BRAND)
    if sub:
        txt(s, sub, 0.7, 0.78, 11.5, 0.36,
            size=13, color=MUTED, italic=True)
    rect(s, 0.5, 1.1, 12.83, 0.04, BORDER)


def pgnum(s, n):
    txt(s, f"{n} / {TOTAL}", 12.0, 7.15, 1.2, 0.3,
        size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def b(text, indent=0):
    return "    " * indent + "  ▸  " + text


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, 13.33, 7.5, SIDEBAR)
rect(s, 0, 0, 13.33, 0.07, BRAND)
rect(s, 0, 7.43, 13.33, 0.07, BRAND)

# decorative glow circles
for cx, cy, sz, col in [
    (10.2, 0.3, 3.8, SURFACE),
    (9.8,  0.1, 2.2, SURFACE2),
    (0.5,  5.5, 3.0, SURFACE),
]:
    circ = s.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    circ.fill.solid(); circ.fill.fore_color.rgb = col
    circ.line.fill.background()

# brand stripe
rect(s, 0, 0, 0.6, 7.5, RGBColor(0x00, 0x90, 0xAA))

txt(s, "FileVault", 1.0, 1.6, 11.2, 1.4,
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "\U0001f510", 5.8, 0.9, 1.8, 1.0,
    size=42, align=PP_ALIGN.CENTER, color=BRAND)

rect(s, 3.2, 3.15, 7.0, 0.04, BRAND)

txt(s, "Enterprise Secure File Management & Integrity Monitoring",
    1.0, 3.3, 11.2, 0.6, size=18, color=BRAND, align=PP_ALIGN.CENTER)
txt(s, "Cybersecurity Student Reference  ·  Full Technical Breakdown",
    1.0, 4.05, 11.2, 0.45, size=13, color=MUTED, align=PP_ALIGN.CENTER,
    italic=True)
txt(s, "Author: Dahirou-Bachar   ·   Python / Flask / PostgreSQL / Railway",
    1.0, 4.9, 11.2, 0.4, size=12, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 3.2, 5.55, 7.0, 0.04, ACCENT)
txt(s, "Stack: Flask · SQLAlchemy · Fernet AES-256 · FIDO2 · TOTP · Watchdog",
    1.0, 5.7, 11.2, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)
pgnum(s, 1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Project Overview", "What is FileVault?")

txt(s, "FileVault is an enterprise-grade secure file vault with real-time integrity monitoring, "
       "multi-factor authentication, role-based clearance levels, and end-to-end AES-256 encryption.",
    0.7, 1.2, 12.1, 0.65, size=14, color=TEXT)

cards = [
    ("\U0001f512  Encrypted Storage",   "AES-256 Fernet per-file\nkey derivation via HKDF",   BRAND),
    ("\U0001f6e1  RBAC + Clearance",    "5-level classification\nL1 TOP SECRET → L5 UNCLASSIFIED", INDIGO),
    ("\U0001f511  MFA / WebAuthn",      "TOTP + FIDO2 biometric\nPhishing-resistant login",    VIOLET),
    ("\U0001f4ca  File Integrity (FIM)","Real-time watchdog +\nScheduled SHA-256 checks",     ACCENT),
    ("\U0001f4cb  Audit Logging",       "Every action logged with\nIP, user, timestamp",       WARNING),
    ("✅  Approval Workflow",       "Admin 3-factor approval\nbefore sensitive operations", SUCCESS),
]
for i, (title, body, col) in enumerate(cards):
    c = i % 3; r = i // 3
    lft = 0.6 + c * 4.24; tp = 2.05 + r * 2.3
    rect(s, lft, tp, 3.95, 2.1, SURFACE)
    rect(s, lft, tp, 3.95, 0.06, col)
    txt(s, title, lft+0.15, tp+0.15, 3.65, 0.45, size=13, bold=True, color=col)
    txt(s, body,  lft+0.15, tp+0.68, 3.65, 1.2,  size=12, color=TEXT)
pgnum(s, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Technology Stack
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Technology Stack", "Libraries, frameworks and services")

stack = [
    ("Flask 3.0.3",            "Web framework – routing, sessions, middleware",           BRAND),
    ("SQLAlchemy 2.0 + Alembic","ORM + database migrations (PostgreSQL)",                     BRAND),
    ("cryptography 43.0.3",    "Fernet AES-256-CBC + HMAC + HKDF",                           ACCENT),
    ("Flask-Login",            "Session & authentication management",                          INDIGO),
    ("Flask-Limiter 4.1.1",   "Rate limiting – brute-force protection",                 INDIGO),
    ("Flask-WTF",              "CSRF token protection on all forms",                           INDIGO),
    ("webauthn 2.2.0",         "FIDO2 / WebAuthn biometric authentication",                   VIOLET),
    ("pyotp 2.9.0",            "TOTP MFA (Google Authenticator compatible)",                   VIOLET),
    ("Flask-SocketIO 5.3.6",  "Real-time push alerts to admin dashboard",                     WARNING),
    ("Watchdog 4.0.1",         "Real-time filesystem event monitoring",                        WARNING),
    ("Resend",                 "Transactional email API (alerts, verification)",               MUTED),
    ("Gunicorn + Railway",     "Production WSGI server & cloud deployment",                    SUCCESS),
]
for i, (lib, desc, col) in enumerate(stack):
    c = i % 2; r = i // 2
    lft = 0.6 + c * 6.35; tp = 1.22 + r * 0.52
    bg = SURFACE if r % 2 == 0 else SURFACE2
    rect(s, lft, tp, 6.05, 0.47, bg)
    rect(s, lft, tp, 0.05, 0.47, col)
    txt(s, lib,  lft+0.15, tp+0.07, 2.2,  0.33, size=11, bold=True, color=col)
    txt(s, desc, lft+2.4,  tp+0.07, 3.55, 0.33, size=11, color=TEXT)
pgnum(s, 3)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Authentication
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Authentication & Session Security")

mtxt(s, [
    ("STRENGTHS", True, BRAND),
    b("Password hashing – Werkzeug PBKDF2 (not MD5 / plain SHA1)"),
    b("Policy enforced – 8+ chars, uppercase, number, symbol"),
    b("TOTP MFA – time-window verified; Google Authenticator compatible"),
    b("WebAuthn/FIDO2 – fingerprint / Face ID; sign-count blocks cloned authenticators"),
    b("Rate limiting – login: 10/min · 50/hr  |  register: 5/min · 20/hr"),
    b("Account lockout – 5 failures → 15-minute lockout"),
    b("Session cookies – HttpOnly · SameSite=Lax · no remember-me"),
    b("Email verification – required before account activation"),
], 0.7, 1.2, 7.3, 5.5, size=13)

rect(s, 8.3, 1.2, 4.7, 5.9, SURFACE)
rect(s, 8.3, 1.2, 4.7, 0.06, ACCENT)
mtxt(s, [
    ("WEAKNESSES", True, ACCENT),
    b("No idle session timeout—session lives until\nbrowser closes"),
    b("SESSION_COOKIE_SECURE defaults\nto False; must be set via env var"),
    b("Password reset requires TOTP—users\nwithout MFA cannot self-serve"),
    b("Login reveals if email exists\n(email enumeration possible)"),
    b("Account lockout is in-memory—\nresets on every app restart"),
], 8.45, 1.3, 4.45, 5.7, size=12)
pgnum(s, 4)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — Encryption
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Encryption & Data Protection")

steps = [
    ("UPLOAD",   "File received\nby Flask"),
    ("KEY GEN",  "HKDF-SHA256\nderives unique\nper-file key"),
    ("ENCRYPT",  "Fernet\nAES-256-CBC\n+ HMAC-SHA256"),
    ("STORE",    "UUID-named\nciphertext saved\nto vault_uploads/"),
    ("DOWNLOAD", "Password\nconfirmed →\ndecrypt in memory"),
]
for i, (label, body) in enumerate(steps):
    lft = 0.5 + i * 2.55
    col = [BRAND, INDIGO, VIOLET, ACCENT, SUCCESS][i]
    rect(s, lft, 1.22, 2.2, 1.9, SURFACE)
    rect(s, lft, 1.22, 2.2, 0.06, col)
    txt(s, label, lft+0.12, 1.28, 1.96, 0.38, size=11, bold=True, color=col)
    txt(s, body,  lft+0.12, 1.72, 1.96, 1.3,  size=11, color=TEXT)
    if i < 4:
        txt(s, "→", lft+2.2, 1.9, 0.32, 0.45,
            size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 0.6, 3.35, 8.3, 3.75, SURFACE)
rect(s, 0.6, 3.35, 8.3, 0.06, BRAND)
mtxt(s, [
    ("KEY FACTS", True, BRAND),
    b("Fernet = AES-256-CBC + HMAC-SHA256 (authenticated encryption)"),
    b("HKDF derives a different key for every single file (per-file uniqueness)"),
    b("Master key stored as ENCRYPTION_KEY environment variable (base64)"),
    b("Decryption gate: user must re-enter account password on every download"),
    b("Files stored as UUID blobs – original filename never appears on disk"),
], 0.75, 3.45, 8.05, 3.55, size=13)

rect(s, 9.1, 3.35, 3.9, 3.75, SURFACE)
rect(s, 9.1, 3.35, 3.9, 0.06, WARNING)
mtxt(s, [
    ("RISKS", True, WARNING),
    b("No key rotation mechanism"),
    b("No HSM integration"),
    b("HKDF info = stored filename\n(deterministic, not random salt)"),
    b("Full file decrypted into RAM\n(risk on large files)"),
    b("Plaintext metadata in DB\n(size, MIME, original name)"),
], 9.25, 3.45, 3.65, 3.55, size=11)
pgnum(s, 5)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — RBAC & Access Control
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Role-Based Access Control (RBAC)", "Clearance levels & approval workflow")

levels = [
    ("L1", "TOP SECRET",    ACCENT,  0.45),
    ("L2", "SECRET",        VIOLET,  1.35),
    ("L3", "CONFIDENTIAL",  INDIGO,  2.25),
    ("L4", "RESTRICTED",    BRAND,   3.15),
    ("L5", "UNCLASSIFIED",  SUCCESS, 4.05),
]
for lvl, name, col, tp in levels:
    w = 5.2 - levels.index((lvl, name, col, tp)) * 0.4
    lft = 0.5 + (5.2 - w) / 2
    rect(s, lft, tp + 1.18, w, 0.72, SURFACE)
    rect(s, lft, tp + 1.18, 0.06, 0.72, col)
    txt(s, f"{lvl}   {name}", lft+0.2, tp+1.26, w-0.3, 0.52,
        size=13, bold=True, color=col)

txt(s, "3-FACTOR APPROVAL WORKFLOW", 6.7, 1.2, 6.4, 0.4,
    size=14, bold=True, color=BRAND)
workflow = [
    ("1", "User submits request + justification",      BRAND),
    ("2", "Admin reviews and approves / rejects",       INDIGO),
    ("3", "One-time execution token issued (30 min)",   VIOLET),
    ("4", "User confirms with account password",        WARNING),
    ("5", "Action executes; token consumed – no replay", SUCCESS),
]
for i, (step, desc, col) in enumerate(workflow):
    tp = 1.75 + i * 0.92
    rect(s, 6.7,  tp, 0.52, 0.68, col)
    txt(s, step, 6.7, tp+0.1, 0.52, 0.48,
        size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(s, 7.26, tp, 5.9, 0.68, SURFACE)
    txt(s, desc, 7.4, tp+0.12, 5.65, 0.46, size=12, color=TEXT)
pgnum(s, 6)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — FIM
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "File Integrity Monitoring (FIM)", "Three-layer real-time detection system")

monitors = [
    ("⏱  Periodic Scanner",  "Every 5 minutes\nSHA-256 hash vs baseline\nCovers all due-check files",       BRAND),
    ("\U0001f504  Recovery Scanner","Every 15 minutes\nRe-checks files with open\nalerts for resolution",       INDIGO),
    ("\U0001f441  Watchdog",        "Real-time OS filesystem events\nDetects changes the instant\nthey happen", ACCENT),
]
for i, (title, body, col) in enumerate(monitors):
    lft = 0.6 + i * 4.25
    rect(s, lft, 1.22, 3.95, 2.4, SURFACE)
    rect(s, lft, 1.22, 3.95, 0.07, col)
    txt(s, title, lft+0.15, 1.3,  3.65, 0.5,  size=14, bold=True, color=col)
    txt(s, body,  lft+0.15, 1.9,  3.65, 1.6,  size=12, color=TEXT)

txt(s, "SEVERITY SCORING – 8 Factors", 0.7, 3.88, 12.0, 0.38,
    size=14, bold=True, color=BRAND)

factors = [
    ("+4", "Critical file type (.exe .sh .py)",   ACCENT),
    ("+4", "Double extension attack (.pdf.exe)",   ACCENT),
    ("+3", "File missing from disk",               VIOLET),
    ("+3", "Repeat offender (3+ prior alerts)",    VIOLET),
    ("+2", "File tampered (hash mismatch)",         INDIGO),
    ("+2", "Bulk event (3+ alerts in 5 min)",       INDIGO),
    ("+2", "Size delta > 50%",                      BRAND),
    ("+1", "Watchdog-detected (real-time)",         SUCCESS),
]
for i, (score, desc, col) in enumerate(factors):
    c = i % 4; r = i // 4
    lft = 0.6 + c * 3.2; tp = 4.42 + r * 0.74
    rect(s, lft,       tp, 0.55, 0.63, col)
    txt(s, score, lft, tp+0.08, 0.55, 0.46,
        size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(s, lft+0.58, tp, 2.55, 0.63, SURFACE)
    txt(s, desc, lft+0.7, tp+0.1, 2.38, 0.44, size=11, color=TEXT)

txt(s, "8+ = CRITICAL   5–7 = HIGH   3–4 = MEDIUM   1–2 = LOW   0 = INFO",
    0.7, 6.88, 12.0, 0.42, size=12, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
pgnum(s, 7)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Input Validation
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Input Validation & Injection Defences")

defences = [
    ("SQL Injection",     SUCCESS,
     "SQLAlchemy ORM – parameterised queries by default.\nilike() for search; no raw SQL with user input."),
    ("XSS",              WARNING,
     "Jinja2 auto-escapes {{ }} by default.\nCSP header set on every response.\nWeak: 'unsafe-inline' still in CSP."),
    ("CSRF",             SUCCESS,
     "Flask-WTF CSRFProtect – tokens embedded in\nevery form. All state-changing routes protected."),
    ("Path Traversal",   SUCCESS,
     "secure_filename() on all uploads.\nFiles stored as UUID blobs – original path discarded."),
    ("File Upload Abuse",WARNING,
     "Extension whitelist (~20 types), 50 MB limit.\nMIME type NOT validated against whitelist – weakness."),
    ("Command Injection",SUCCESS,
     "No shell commands executed anywhere.\nos.path.join() for all path construction."),
]
for i, (attack, status, detail) in enumerate(defences):
    c = i % 2; r = i // 2
    lft = 0.6 + c * 6.35; tp = 1.22 + r * 2.05
    rect(s, lft, tp, 6.05, 1.92, SURFACE)
    rect(s, lft, tp, 0.07, 1.92, status)
    icon = "✔" if status == SUCCESS else "⚠"
    txt(s, icon,   lft+0.15, tp+0.1,  0.4,  0.42, size=14, bold=True, color=status)
    txt(s, attack, lft+0.6,  tp+0.1,  5.3,  0.42, size=13, bold=True, color=status)
    txt(s, detail, lft+0.2,  tp+0.6,  5.7,  1.2,  size=11, color=TEXT)
pgnum(s, 8)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Security Headers
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Security Headers & Web Configuration")

headers_data = [
    ("X-Content-Type-Options: nosniff",                     "Prevents MIME-sniffing attacks",                SUCCESS),
    ("X-Frame-Options: SAMEORIGIN",                          "Clickjacking protection (same-origin only)",    SUCCESS),
    ("Referrer-Policy: strict-origin-when-cross-origin",     "Limits referrer information leakage",           SUCCESS),
    ("Permissions-Policy: geolocation=(), mic=(), camera=()", "Disables dangerous browser APIs",              SUCCESS),
    ("X-XSS-Protection: 1; mode=block",                      "Legacy header – deprecated by modern browsers",  WARNING),
    ("CSP: default-src 'self' ... 'unsafe-inline'",          "Inline scripts allowed – weakens XSS protection", ACCENT),
    ("HSTS header",                                          "MISSING – browser not forced to HTTPS",    DANGER),
    ("SESSION_COOKIE_SECURE",                                "Defaults to False – must be set via env var", DANGER),
    ("SocketIO cors_allowed_origins='*'",                    "Open to all origins – should be locked to domain", DANGER),
]
for i, (h, note, status) in enumerate(headers_data):
    tp = 1.22 + i * 0.57
    icon = "✔" if status == SUCCESS else ("⚠" if status == WARNING else "✘")
    bg = SURFACE if i % 2 == 0 else SURFACE2
    rect(s, 0.5, tp, 12.5, 0.52, bg)
    rect(s, 0.5, tp, 0.06, 0.52, status)
    txt(s, icon, 0.58, tp+0.06, 0.38, 0.38, size=13, bold=True,
        color=status, align=PP_ALIGN.CENTER)
    txt(s, h,    1.02, tp+0.06, 7.1,  0.38, size=10, bold=True, color=status)
    txt(s, note, 8.2,  tp+0.06, 4.7,  0.38, size=10, color=MUTED)
pgnum(s, 9)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Audit Logging
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Audit Logging & Activity Monitoring")

txt(s, "Every significant action is recorded with:  User · Action · Detail · IP Address · Timestamp",
    0.7, 1.2, 12.1, 0.48, size=13, color=TEXT)

events = [
    ("login  /  login_fail  /  logout",             BRAND),
    ("register  /  email_verified",                 BRAND),
    ("upload  /  download  /  delete  /  rename",   INDIGO),
    ("request_submitted  /  request_approved",       VIOLET),
    ("tamper_detected  /  baseline_reset",           ACCENT),
    ("mfa_enabled  /  webauthn_registered",          SUCCESS),
    ("admin_action  /  profile_update",              WARNING),
    ("export_csv  /  export_pdf",                   MUTED),
]
for i, (evt, col) in enumerate(events):
    c = i % 2; r = i // 2
    lft = 0.6 + c * 6.35; tp = 1.82 + r * 0.7
    rect(s, lft, tp, 6.05, 0.58, SURFACE)
    rect(s, lft, tp, 0.06, 0.58, col)
    txt(s, evt, lft+0.2, tp+0.1, 5.7, 0.38, size=12, color=col)

rect(s, 0.6, 4.75, 12.05, 2.0, SURFACE)
rect(s, 0.6, 4.75, 12.05, 0.06, WARNING)
mtxt(s, [
    ("LIMITATIONS", True, WARNING),
    b("Logs stored only in the application database – no SIEM / syslog integration"),
    b("No log retention or rotation policy – could grow unbounded"),
    b("All admins can read all logs – no sensitivity segregation"),
    b("No content hash logged on download – cannot verify integrity of what was sent"),
], 0.75, 4.85, 11.8, 1.82, size=12)
pgnum(s, 10)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — Secrets Management
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Secrets Management & Configuration")

secrets_data = [
    ("SECRET_KEY",      "Flask session signing key",        WARNING,
     "Has insecure default string – deployment may miss changing it"),
    ("ENCRYPTION_KEY",  "Master Fernet encryption key",     SUCCESS,
     "No default – app fails to start without it (good)"),
    ("DATABASE_URL",    "PostgreSQL connection string",     SUCCESS,
     "Handles postgres:// → postgresql:// URI fix automatically"),
    ("RESEND_API_KEY",  "Email service API key",            DANGER,
     "Leaked key → attacker can send emails as FileVault"),
    ("ADMIN_EMAIL",     "Bootstrapped admin account",       WARNING,
     "Auto-created on startup via prestart.py"),
    ("WEBAUTHN_RP_ID",  "FIDO2 relying party ID (domain)",  WARNING,
     "Must match production domain exactly or FIDO2 breaks"),
]
rect(s, 0.5, 1.22, 12.5, 0.5, SURFACE2)
for j, (hdr, lft) in enumerate(zip(
    ["Variable", "Purpose", "Note"],
    [0.65, 2.95, 6.55]
)):
    txt(s, hdr, lft, 1.3, [2.2, 3.5, 5.8][j], 0.36,
        size=12, bold=True, color=BRAND)

for i, (var, purpose, risk, note) in enumerate(secrets_data):
    tp = 1.78 + i * 0.72
    bg = SURFACE if i % 2 == 0 else SURFACE2
    rect(s, 0.5, tp, 12.5, 0.67, bg)
    rect(s, 0.5, tp, 0.06, 0.67, risk)
    icon = "✔" if risk == SUCCESS else ("⚠" if risk == WARNING else "✘")
    txt(s, icon,    0.62, tp+0.12, 0.28, 0.42, size=12, bold=True,
        color=risk, align=PP_ALIGN.CENTER)
    txt(s, var,     0.95, tp+0.1,  1.9,  0.45, size=11, bold=True, color=risk)
    txt(s, purpose, 2.95, tp+0.1,  3.45, 0.45, size=11, color=TEXT)
    txt(s, note,    6.55, tp+0.1,  6.3,  0.45, size=10, color=MUTED)
pgnum(s, 11)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Rate Limiting
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Rate Limiting & DoS Protection")

limits = [
    ("Login",                  "10 / min   ·   50 / hr",   BRAND),
    ("Register",               "5 / min    ·   20 / hr",   BRAND),
    ("File Upload",            "20 / hr",                        INDIGO),
    ("Decrypt / Download",     "5 / min",                        VIOLET),
    ("Manual Integrity Check", "30 / min",                       SUCCESS),
    ("File Rename / Delete",   "Per-operation limits",           MUTED),
]
for i, (endpoint, limit, col) in enumerate(limits):
    c = i % 2; r = i // 2
    lft = 0.6 + c * 6.35; tp = 1.22 + r * 0.88
    rect(s, lft, tp, 6.05, 0.76, SURFACE)
    rect(s, lft, tp, 0.06, 0.76, col)
    txt(s, endpoint, lft+0.2,  tp+0.08, 3.5,  0.56, size=12, bold=True, color=col)
    txt(s, limit,    lft+3.75, tp+0.08, 2.15, 0.56, size=12, color=BRAND,
        align=PP_ALIGN.RIGHT)

rect(s, 0.6, 4.0, 12.05, 3.15, SURFACE)
rect(s, 0.6, 4.0, 12.05, 0.06, DANGER)
mtxt(s, [
    ("CRITICAL WEAKNESSES", True, DANGER),
    b("In-memory storage – rate limit counters RESET on every app restart (Railway auto-restarts frequently)"),
    b("Not distributed – counters are NOT shared across multiple app instances / load balancer nodes"),
    b("Account lockout also in-memory – same problem; lockouts lifted by restarts"),
    b("Fix: Redis-backed Flask-Limiter for persistent, shared counters across all instances"),
    b("IP-based limiting relies on X-Forwarded-For – trust depends on correct ProxyFix configuration"),
], 0.75, 4.1, 11.8, 2.97, size=12)
pgnum(s, 12)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Password Reset
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Password Reset & Account Recovery")

steps = [
    ("1  Request",      "User enters email\non Forgot Password",    BRAND),
    ("2  TOTP Gate",    "Must verify TOTP\nbefore proceeding",      INDIGO),
    ("3  Token Gen",    "secrets.token_urlsafe(32)\n256-bit random; 1 hr TTL", VIOLET),
    ("4  Email Link",   "Secure link sent\nvia Resend API",         ACCENT),
    ("5  New Password", "Token consumed\n(single-use only)",        SUCCESS),
]
for i, (label, body, col) in enumerate(steps):
    lft = 0.5 + i * 2.55
    rect(s, lft, 1.22, 2.2, 1.9, SURFACE)
    rect(s, lft, 1.22, 2.2, 0.06, col)
    txt(s, label, lft+0.12, 1.28, 1.96, 0.42, size=12, bold=True, color=col)
    txt(s, body,  lft+0.12, 1.78, 1.96, 1.22, size=11, color=TEXT)
    if i < 4:
        txt(s, "→", lft+2.2, 1.95, 0.32, 0.42,
            size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

rect(s, 0.6, 3.35, 5.9, 3.8, SURFACE)
rect(s, 0.6, 3.35, 5.9, 0.06, SUCCESS)
mtxt(s, [
    ("STRENGTHS", True, SUCCESS),
    b("256-bit cryptographically random token"),
    b("1-hour expiry – short attack window"),
    b("Single-use – consumed on first click"),
    b("TOTP gate prevents email-only account takeover"),
    b("Email must be verified first – prevents\npre-registered hijacking"),
], 0.75, 3.45, 5.65, 3.6, size=12)

rect(s, 6.75, 3.35, 6.1, 3.8, SURFACE)
rect(s, 6.75, 3.35, 6.1, 0.06, ACCENT)
mtxt(s, [
    ("WEAKNESSES", True, ACCENT),
    b("Users WITHOUT MFA cannot reset their own\npassword – must contact admin"),
    b("No backup codes or recovery email fallback"),
    b("No rate limiting on reset email generation\n– can spam the email API"),
    b("Token comparison could use hmac.compare_\ndigest to prevent timing attacks"),
], 6.9, 3.45, 5.85, 3.6, size=12)
pgnum(s, 13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — Risk Matrix
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Risk Matrix & Threat Summary")

risks = [
    ("Master Encryption Key Leak",        "CRITICAL", DANGER,
     "All files decryptable; no key rotation or HSM"),
    ("Session Hijacking",                 "HIGH",     ACCENT,
     "HTTPOnly on, but HTTPS not enforced by default"),
    ("Admin Account Compromise",          "HIGH",     ACCENT,
     "Full system takeover; force MFA on admin accounts"),
    ("Baseline Reset by Malicious Admin", "MEDIUM",   WARNING,
     "Can legitimise tampering; needs second admin approval"),
    ("TOCTOU Race in FIM Hashing",        "MEDIUM",   WARNING,
     "File can be modified during multi-chunk SHA-256"),
    ("Rate Limit Resets on Restart",      "MEDIUM",   WARNING,
     "In-memory counters cleared on every Railway restart"),
    ("CSP unsafe-inline",                 "MEDIUM",   INDIGO,
     "Inline scripts allowed; weakens XSS defence"),
    ("MIME Type Not Validated on Upload", "MEDIUM",   INDIGO,
     "Attacker renames .exe to .txt and uploads it"),
    ("Email Enumeration on Login",        "LOW",      SUCCESS,
     "Error message reveals whether email exists in system"),
    ("Log Retention / SIEM Gap",          "LOW",      SUCCESS,
     "No log rotation; no external SIEM integration"),
    ("Symlink Swap Attack",               "LOW",      SUCCESS,
     "Watchdog may hash symlink target instead of link"),
]
rect(s, 0.5, 1.22, 12.5, 0.48, SURFACE2)
for j, (hdr, lft) in enumerate(zip(
    ["Threat", "Severity", "Notes"],
    [0.65, 7.1, 8.95]
)):
    txt(s, hdr, lft, 1.3, [6.3, 1.7, 3.9][j], 0.34, size=12, bold=True, color=BRAND)

for i, (threat, severity, col, note) in enumerate(risks):
    tp = 1.76 + i * 0.5
    bg = SURFACE if i % 2 == 0 else SURFACE2
    rect(s, 0.5, tp, 12.5, 0.46, bg)
    rect(s, 0.5, tp, 0.06, 0.46, col)
    txt(s, threat,   0.65, tp+0.06, 6.3,  0.34, size=10, color=TEXT)
    txt(s, severity, 7.1,  tp+0.06, 1.7,  0.34, size=10, bold=True, color=col)
    txt(s, note,     8.95, tp+0.06, 3.9,  0.34, size=9,  color=MUTED)
pgnum(s, 14)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Recommendations
# ═══════════════════════════════════════════════════════════════════
s = slide()
header(s, "Recommendations", "Prioritised security improvements")

recs = [
    ("\U0001f534  IMMEDIATE",  DANGER, [
        "Enforce HTTPS + add HSTS header",
        "Set SESSION_COOKIE_SECURE = True by default",
        "Remove 'unsafe-inline' from CSP; use nonces",
        "Add idle session timeout (15–30 min) with re-auth",
    ]),
    ("\U0001f7e1  SHORT-TERM", WARNING, [
        "Backup codes / recovery email for users without MFA",
        "Redis-backed rate limiting (persistent across restarts)",
        "Require admin MFA confirmation before baseline reset",
        "Validate MIME type against extension whitelist on upload",
    ]),
    ("\U0001f7e2  MEDIUM-TERM", SUCCESS, [
        "Key rotation mechanism + HSM for ENCRYPTION_KEY",
        "Random per-file salt in HKDF (store salt in DB)",
        "Restrict SocketIO CORS to app domain only",
        "SIEM / syslog integration for audit logs",
    ]),
]
for i, (title, col, items) in enumerate(recs):
    lft = 0.6 + i * 4.24
    rect(s, lft, 1.22, 3.95, 5.95, SURFACE)
    rect(s, lft, 1.22, 3.95, 0.07, col)
    txt(s, title, lft+0.15, 1.3, 3.65, 0.45, size=13, bold=True, color=col)
    lines = [("", False, col)]
    for item in items:
        lines.append((b(item), False, TEXT))
    mtxt(s, lines, lft+0.15, 1.9, 3.68, 5.1, size=12)
pgnum(s, 15)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Key Takeaways
# ═══════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, 13.33, 7.5, SIDEBAR)
rect(s, 0, 0, 13.33, 0.07, BRAND)
rect(s, 0, 7.43, 13.33, 0.07, BRAND)
rect(s, 0, 0, 0.6,  7.5,  RGBColor(0x00, 0x90, 0xAA))

txt(s, "Key Takeaways", 0.8, 0.18, 11.8, 0.72,
    size=34, bold=True, color=BRAND, align=PP_ALIGN.CENTER)
rect(s, 3.2, 0.95, 7.0, 0.05, ACCENT)

takeaways = [
    ("Defence-in-Depth",      BRAND,
     "Download = clearance check + admin approval + one-time token + password confirm. Four independent gates."),
    ("FIM Severity Scoring",  INDIGO,
     "Not binary tampered/OK. 8-factor weighted system gives nuanced threat intelligence to admins."),
    ("Per-File Key Derivation", VIOLET,
     "HKDF gives each file a unique encryption key. Compromising one file does NOT expose others."),
    ("WebAuthn / FIDO2",      ACCENT,
     "Phishing-resistant login. Credentials are cryptographically bound to the origin domain – cannot be phished."),
    ("TOTP-Gated Password Reset", SUCCESS,
     "Password reset requires the TOTP app. Email hijack alone is not enough to take over an account."),
]
for i, (title, col, body) in enumerate(takeaways):
    tp = 1.12 + i * 1.2
    rect(s, 0.75, tp, 0.07, 1.05, col)
    rect(s, 0.88, tp, 12.0, 1.05, SURFACE)
    txt(s, title, 1.05, tp+0.08, 4.0,  0.38, size=14, bold=True, color=col)
    txt(s, body,  1.05, tp+0.52, 11.5, 0.46, size=12, color=TEXT)

txt(s, "FileVault  ·  Built by Dahirou-Bachar  ·  Cybersecurity Final Project",
    0.8, 7.1, 12.0, 0.35, size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
pgnum(s, 16)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out = r"c:\Users\user\Desktop\FINAL PROJECT\FileVault_Cybersecurity_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
